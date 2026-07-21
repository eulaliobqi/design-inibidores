"""Generate 22-slide PowerPoint presentation via python-pptx."""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

UFV_BLUE = RGBColor(0, 51, 102)  # #003366
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)

def load_slides_spec():
    """Load 22 slides from JSON."""
    return json.loads(Path("scripts/canva_slides_data.json").read_text(encoding="utf-8"))

def create_blank_slide_layout(prs):
    """Get blank slide layout from presentation."""
    return prs.slide_layouts[6]  # Blank layout

def add_title_slide(prs, title, subtitle):
    """Add title slide (Slide 1)."""
    slide = prs.slides.add_slide(create_blank_slide_layout(prs))

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9.0), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = UFV_BLUE

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9.0), Inches(1.0))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY

def add_content_slide(prs, title, bullets):
    """Add content slide with title + bullets + placeholder for image."""
    slide = prs.slides.add_slide(create_blank_slide_layout(prs))

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = UFV_BLUE

    # Add bullets
    if bullets:
        bullets_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(5.8), Inches(5.0))
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = bullets_frame.paragraphs[0]
            else:
                p = bullets_frame.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.level = 0
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    # Add placeholder for image (right side)
    # Image will be inserted in Task 2 if needed
    image_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(6.8), Inches(1.1), Inches(2.7), Inches(5.0)
    )
    image_placeholder.fill.solid()
    image_placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
    image_placeholder.line.color.rgb = RGBColor(200, 200, 200)

def insert_figures(prs):
    """Insert 4 PNG figures into slides 12, 13, 14, 18 and remove placeholders."""
    figures_map = {
        12: "outputs/figuras_artigo/fig4_fingerprint.png",
        13: "outputs/figuras_artigo/fig1_replicas_md.png",
        14: "outputs/figuras_artigo/fig2_ocupancia_s1.png",
        18: "outputs/figuras_artigo/fig3_cross_species.png",
    }

    for slide_idx, (slide_num, fig_path) in enumerate(figures_map.items()):
        # slide_idx 0-indexed, slide_num 1-indexed
        # Slide 12 is at prs.slides[11], etc.
        slide = prs.slides[slide_num - 1]
        fig_full_path = Path(fig_path)

        if not fig_full_path.exists():
            print(f"  WARNING: {fig_path} not found, skipping Slide {slide_num}")
            continue

        # Remove placeholder rectangle (last shape, typically)
        # Iterate backwards to avoid index shifts
        for shape_idx in range(len(slide.shapes) - 1, -1, -1):
            shape = slide.shapes[shape_idx]
            # Check if it's the placeholder rectangle (has fill, no text)
            if hasattr(shape, "fill") and not hasattr(shape, "text_frame"):
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    print(f"  Removed placeholder from Slide {slide_num}")
                    break
                except Exception as e:
                    print(f"  Could not remove placeholder: {e}")

        # Insert figure
        try:
            slide.shapes.add_picture(
                str(fig_full_path),
                Inches(6.8), Inches(1.1),
                width=Inches(2.7)
            )
            print(f"  Inserted {fig_path} into Slide {slide_num}")
        except Exception as e:
            print(f"  ERROR inserting figure into Slide {slide_num}: {e}")

def main():
    """Build all 22 slides."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides_data = load_slides_spec()

    for i, slide_info in enumerate(slides_data['slides']):
        slide_num = slide_info['slide_num']
        title = slide_info['title']
        subtitle = slide_info.get('subtitle')
        bullets = slide_info.get('bullets', [])

        print(f"Adding Slide {slide_num}: {title[:50]}...")

        if slide_num == 1:
            add_title_slide(prs, title, subtitle)
        else:
            add_content_slide(prs, title, bullets)

    # Insert figures into specific slides
    print("\nInserting figures...")
    insert_figures(prs)

    # Save
    output_path = Path("apresentacao_design_inibidores.pptx")
    prs.save(str(output_path))
    print(f"\nSalvo: {output_path.resolve()}")
    print(f"Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
