"""Apresentação de slides: Design Racional de Inibidores Peptídicos via IA Generativa.

Uso:
    python criar_slides.py
    # Gera: apresentacao_design_inibidores.pptx

Atualizado 2026-06-27 — Fase 1-3 concluída:
  RFdiffusion 330 | ProteinMPNN 24.513 | Vina 880 | PyRosetta 10 | MD 5/5 | ML RF | Specificity 20/20 | Cleavage 0/20
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ─── Paleta de cores ──────────────────────────────────────────────────────────
C_AZUL_ESCURO  = RGBColor(0x0D, 0x2B, 0x55)
C_AZUL_MEDIO   = RGBColor(0x1A, 0x4F, 0x8A)
C_CIANO        = RGBColor(0x00, 0xB4, 0xD8)
C_VERDE        = RGBColor(0x06, 0xD6, 0xA0)
C_AMARELO      = RGBColor(0xFF, 0xD1, 0x66)
C_VERMELHO     = RGBColor(0xEF, 0x47, 0x6F)
C_BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
C_CINZA_CLARO  = RGBColor(0xF0, 0xF4, 0xF8)
C_CINZA_TEXTO  = RGBColor(0x33, 0x33, 0x44)
C_LARANJA      = RGBColor(0xFF, 0x7F, 0x11)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

# ─── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line_color=None, line_w=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, italic=False,
             color=C_BRANCO, align=PP_ALIGN.LEFT,
             word_wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box

def bg_dark(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, C_AZUL_ESCURO)

def bg_light(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, C_CINZA_CLARO)
    add_rect(slide, 0, 0, 13.33, 1.1, C_AZUL_ESCURO)

TOTAL_SLIDES = 18

def slide_num(slide, n):
    add_text(slide, f"{n}/{TOTAL_SLIDES}", 12.6, 7.1, 0.7, 0.3,
             font_size=11, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.RIGHT)

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.1, C_AZUL_ESCURO)
    add_rect(slide, 0, 1.05, 13.33, 0.05, C_CIANO)
    add_text(slide, title, 0.4, 0.05, 12.0, 0.65,
             font_size=26, bold=True, color=C_BRANCO)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.68, 12.0, 0.38,
                 font_size=14, color=C_CIANO)

def tag(slide, text, x, y, w=1.8, h=0.35, bg=C_AZUL_MEDIO, fg=C_BRANCO, font_size=13):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, text, x + 0.05, y + 0.02, w - 0.1, h,
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Capa
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_dark(slide)
add_rect(slide, 0, 2.5, 13.33, 0.08, C_CIANO)
add_rect(slide, 0, 4.95, 13.33, 0.05, C_AZUL_MEDIO)

add_text(slide,
    "Design Racional de Inibidores Peptídicos",
    0.6, 0.9, 12.0, 0.9,
    font_size=38, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
add_text(slide,
    "de Tripsinas de Lepidoptera por IA Generativa",
    0.6, 1.7, 12.0, 0.75,
    font_size=32, bold=True, color=C_CIANO, align=PP_ALIGN.CENTER)

add_text(slide,
    "Pipeline Multiagente: RFdiffusion → ProteinMPNN → PyRosetta → Vina → GROMACS",
    0.6, 2.7, 12.0, 0.5,
    font_size=16, italic=True, color=C_CINZA_CLARO, align=PP_ALIGN.CENTER)

add_text(slide,
    "Eulalio B. Q. Santos  |  UFV Viçosa-MG  |  eulalio.santos@ufv.br",
    0.6, 5.15, 12.0, 0.45,
    font_size=15, color=C_CINZA_CLARO, align=PP_ALIGN.CENTER)
add_text(slide,
    "github.com/eulaliobqi/design-inibidores  |  Junho 2026",
    0.6, 5.6, 12.0, 0.4,
    font_size=13, color=C_CIANO, align=PP_ALIGN.CENTER)

techs = ["RFdiffusion", "ProteinMPNN", "PyRosetta 2026", "AutoDock Vina", "GROMACS", "Python 3.10"]
colors = [C_AZUL_MEDIO, C_AZUL_MEDIO, C_VERDE, C_VERDE, C_VERDE, C_AZUL_MEDIO]
for i, (t, c) in enumerate(zip(techs, colors)):
    tag(slide, t, 0.6 + i * 2.05, 6.2, w=1.9, h=0.38, bg=c, font_size=12)

slide_num(slide, 1)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Contextualização
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Contextualização", "Por que inibidores de tripsinas de Lepidoptera?")

add_rect(slide, 0.4, 1.25, 5.8, 5.5, C_AZUL_ESCURO)
add_text(slide, "O PROBLEMA", 0.6, 1.35, 5.4, 0.45,
         font_size=16, bold=True, color=C_CIANO)
itens_prob = [
    "  Spodoptera frugiperda e Agrotis gemmatalis",
    "      causam bilhões USD/ano em perdas agrícolas",
    "",
    "  Inseticidas químicos: resistência crescente,",
    "      toxicidade ambiental e resíduos em alimentos",
    "",
    "  Tripsinas digestivas são essenciais para o",
    "      desenvolvimento larval — alvos comprovados",
    "",
    "  Quais peptídeos inibem essas tripsinas",
    "      com alta seletividade e baixa toxicidade?",
]
for i, item in enumerate(itens_prob):
    add_text(slide, item, 0.55, 1.85 + i * 0.46, 5.5, 0.48,
             font_size=13, color=C_CINZA_CLARO)

add_rect(slide, 6.6, 1.25, 6.3, 5.5, C_BRANCO)
add_rect(slide, 6.6, 1.25, 6.3, 0.05, C_CIANO)
add_text(slide, "A OPORTUNIDADE DA IA GENERATIVA", 6.8, 1.35, 5.9, 0.45,
         font_size=15, bold=True, color=C_AZUL_ESCURO)
itens_oport = [
    "RFdiffusion (Nature 2023): gera backbones",
    "    proteicos condicionados a sítios de ligação",
    "",
    "ProteinMPNN (Science 2022): otimiza",
    "    sequências para backbones de interesse",
    "",
    "PyRosetta (2026.25): refinamento de interface",
    "    FastRelax + InterfaceAnalyzerMover",
    "",
    "AutoDock Vina: triagem de afinidade em",
    "    larga escala (194 poses reais, kcal/mol)",
    "",
    "Dataset ML/DL: 24.513 seqs × 41 features",
    "    para treinar modelos preditivos",
]
for i, item in enumerate(itens_oport):
    add_text(slide, item, 6.75, 1.85 + i * 0.36, 5.9, 0.38,
             font_size=12, color=C_CINZA_TEXTO)

slide_num(slide, 2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Objetivos
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Objetivos", "Design computacional racional + dataset para ML/DL")

add_text(slide, "Objetivo geral:", 0.5, 1.25, 12.0, 0.45,
         font_size=19, bold=True, color=C_AZUL_ESCURO)
add_text(slide,
    "Gerar e validar computacionalmente peptídeos inibidores de tripsinas digestivas de "
    "Lepidoptera-praga usando pipeline multiagente de IA generativa estrutural, produzindo "
    "simultaneamente um dataset ML/DL supervisionado com scores Vina + Rosetta como labels.",
    0.5, 1.65, 12.3, 0.9,
    font_size=15, color=C_CINZA_TEXTO)

add_text(slide, "Objetivos específicos:", 0.5, 2.65, 12.0, 0.4,
         font_size=17, bold=True, color=C_AZUL_ESCURO)

esp = [
    ("1.", "Identificar resíduos catalíticos e bolso S1 de 4 tripsinas lepidópteras (ACR157/QCL936/XP273/XP352)"),
    ("2.", "Gerar 330 backbones peptídicos de 20 aa via RFdiffusion ancorados ao sítio S1 de ACR157"),
    ("3.", "Produzir 24.513 sequências binder únicas por ProteinMPNN (extração cadeia B, bug FASTA corrigido)"),
    ("4.", "Anotar 41 features físico-químicas para cada sequência (dataset ML/DL supervisionado)"),
    ("5.", "Obter energias de afinidade Vina reais como labels — 194/194 poses válidas (binders RFdiffusion)"),
    ("6.", "Refinar top-10 candidatos com PyRosetta FastRelax + InterfaceAnalyzerMover (I_sc REF2015)"),
    ("7.", "Validar top-5 por dinâmica molecular 10 ns (GROMACS gmx_mpi, AMBER99SB-ILDN, TIP3P, 300 K)"),
]
for i, (num, texto) in enumerate(esp):
    add_rect(slide, 0.5, 3.1 + i * 0.55, 0.45, 0.42, C_AZUL_MEDIO)
    add_text(slide, num, 0.5, 3.1 + i * 0.55, 0.45, 0.42,
             font_size=15, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, texto, 1.05, 3.13 + i * 0.55, 11.7, 0.42,
             font_size=13, color=C_CINZA_TEXTO)

slide_num(slide, 3)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Pipeline: arquitetura geral
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Pipeline Multiagente", "10 módulos especializados — TODAS as etapas principais concluídas")

steps = [
    ("01", "Structure\nAgent",    "Sítio S1\nHotspots",           C_AZUL_MEDIO, "✓"),
    ("02", "RFdiffusion\nAgent",  "330 backbones\n20 aa reais",   C_AZUL_MEDIO, "✓"),
    ("03", "ProteinMPNN\nAgent",  "24.513\nbinders únicos",       C_AZUL_MEDIO, "✓"),
    ("04", "Rosetta\nAgent",      "PyRosetta\nI_sc REF2015",      C_VERDE,      "✓"),
    ("05", "Docking\nAgent",      "194/194\nVina reais",          C_VERDE,      "✓"),
    ("06", "MD Agent",            "GROMACS\n10 ns",               C_AMARELO,    "◑"),
    ("07", "Ranking\nAgent",      "Score\ncomposto",              C_AZUL_MEDIO, "✓"),
    ("08", "Visualize\nAgent",    "6 figuras\n",                  C_VERDE,      "✓"),
    ("09", "Report\nAgent",       "HTML + MD\n",                  C_VERDE,      "✓"),
    ("10", "Optimize\nAgent",     "Redesign\ntop-10",             C_CINZA_TEXTO,"—"),
]

colors_step = {
    "✓": C_VERDE,
    "◑": C_AMARELO,
    "✗": C_VERMELHO,
    "—": C_CINZA_TEXTO,
}

for i, (num, name, desc, c, status) in enumerate(steps):
    col = i % 5
    row = i // 5
    x = 0.35 + col * 2.52
    y = 1.25 + row * 2.7
    sc = colors_step[status]
    add_rect(slide, x, y, 2.25, 2.3, sc)
    add_rect(slide, x + 0.07, y + 0.07, 2.11, 2.16, C_BRANCO)
    add_text(slide, num, x + 0.1, y + 0.1, 0.5, 0.35,
             font_size=12, bold=True, color=sc)
    add_text(slide, status, x + 1.8, y + 0.1, 0.5, 0.35,
             font_size=16, bold=True, color=sc, align=PP_ALIGN.RIGHT)
    add_text(slide, name, x + 0.1, y + 0.45, 2.05, 0.65,
             font_size=13, bold=True, color=C_AZUL_ESCURO)
    add_text(slide, desc, x + 0.1, y + 1.1, 2.05, 0.85,
             font_size=12, color=C_CINZA_TEXTO)
    if col < 4:
        add_text(slide, "→", x + 2.25, y + 0.9, 0.27, 0.45,
                 font_size=20, bold=True, color=C_CIANO, align=PP_ALIGN.CENTER)

for status, label, c in [("✓","Concluído",C_VERDE),("◑","Em execução",C_AMARELO),("—","Aguarda MD",C_CINZA_TEXTO)]:
    idx = ["✓","◑","—"].index(status)
    add_rect(slide, 0.35 + idx * 2.5, 7.1, 0.3, 0.25, c)
    add_text(slide, f" {label}", 0.65 + idx * 2.5, 7.07, 2.0, 0.3,
             font_size=12, color=C_CINZA_TEXTO)

slide_num(slide, 4)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Alvo molecular
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Alvo Molecular", "Tripsinas digestivas de Lepidoptera-praga")

add_text(slide,
    "4 modelos estruturais de tripsinas: ACR157, QCL936, XP273, XP352",
    0.5, 1.2, 12.0, 0.4, font_size=16, color=C_AZUL_ESCURO, bold=True)

headers = ["Modelo", "His/Tyr (Cat)", "Asp (Tríade)", "Ser/Ala (Nuc)", "S1 (Esp)", "Centro (X,Y,Z) Å"]
rows_tab = [
    ["ACR157",    "His69",   "Asp114",  "Ser211",  "Asp205",  "7.09 / -1.15 / -1.63"],
    ["QCL936",    "His92",   "Asp142",  "Ser247",  "Asp241",  "2.94 / 5.55 / -0.32"],
    ["XP273 ★",  "Tyr83",   "Asp132",  "Ser234",  "Ile229",  "1.67 / 8.29 / -0.34"],
    ["XP352",    "His112",  "Asp166",  "Ala268",  "Asp262",  "-1.26 / 5.60 / -5.25"],
    ["Consenso", "-",        "-",       "-",       "-",       "2.61 / 4.57 / -1.89"],
]
widths = [1.4, 1.6, 1.45, 1.55, 1.3, 2.65]
x0, y0 = 0.4, 1.7
for j, (hdr, w) in enumerate(zip(headers, widths)):
    add_rect(slide, x0 + sum(widths[:j]), y0, w, 0.42, C_AZUL_ESCURO)
    add_text(slide, hdr, x0 + sum(widths[:j]) + 0.05, y0 + 0.02, w - 0.1, 0.38,
             font_size=12, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
for i, row in enumerate(rows_tab):
    bg = C_BRANCO if i % 2 == 0 else C_CINZA_CLARO
    if row[0] == "Consenso":
        bg = C_AZUL_MEDIO
    for j, (cell, w) in enumerate(zip(row, widths)):
        add_rect(slide, x0 + sum(widths[:j]), y0 + 0.42 + i * 0.4, w, 0.4, bg)
        fc = C_BRANCO if row[0] == "Consenso" else C_CINZA_TEXTO
        add_text(slide, cell, x0 + sum(widths[:j]) + 0.05, y0 + 0.44 + i * 0.4, w - 0.1, 0.36,
                 font_size=12, color=fc, align=PP_ALIGN.CENTER)

add_rect(slide, 0.4, 4.9, 12.5, 1.45, RGBColor(0xFF, 0xF3, 0xCD))
add_rect(slide, 0.4, 4.9, 0.08, 1.45, C_AMARELO)
add_text(slide, "★ XP273 — Tripsina Atípica", 0.6, 4.95, 5.0, 0.4,
         font_size=15, bold=True, color=C_AZUL_ESCURO)
add_text(slide,
    "Tyr83 substitui His catalítica canônica | Ile229 no bolso S1 (ao invés de Asp) "
    "-> especificidade hidrofóbica\n"
    "Justifica NAO restringir P1 no pipeline — candidatos acidos (D/E) podem atuar via mecanismo nao-canonico",
    0.6, 5.35, 12.0, 0.98, font_size=13, color=C_CINZA_TEXTO)

add_rect(slide, 0.4, 6.45, 12.5, 0.75, C_AZUL_ESCURO)
add_text(slide,
    "Grid box Vina (consenso):  centro = [2.61, 4.57, -1.89] A  |  "
    "tamanho adaptativo = max(40, len*3.6+8) A  |  exaustividade = 8",
    0.6, 6.55, 12.0, 0.55, font_size=14, color=C_BRANCO, align=PP_ALIGN.CENTER)

slide_num(slide, 5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — RFdiffusion: 330 backbones reais
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Geração de Backbones Peptídicos — RFdiffusion REAL",
           "330 backbones de novo | Watson et al., Nature 2023 | checkpoint Complex_base_ckpt.pt 462 MB")

add_rect(slide, 0.4, 1.25, 5.8, 5.8, C_AZUL_ESCURO)
add_text(slide, "RFdiffusion", 0.6, 1.35, 5.4, 0.45,
         font_size=18, bold=True, color=C_CIANO)
add_text(slide, "Instalado em ~/RFdiffusion | Python 3.10 (protein_design_env)", 0.6, 1.75, 5.4, 0.3,
         font_size=11, italic=True, color=C_CINZA_CLARO)

items_rfd = [
    "  ~/RFdiffusion instalado e funcional",
    "  Complex_base_ckpt.pt (462 MB) baixado",
    "  330 backbones gerados para ACR157",
    "",
    "Parâmetros usados:",
    "  noise_scale_ca  = 0.2",
    "  noise_scale_frame = 0.1",
    "  contig: A1-231/0 20-20",
    "  inference.ckpt_override_path=.../models/",
    "",
    "Pre-processamento obrigatório:",
    "  prep_pdbs.py: chain blank -> A",
    "  renumeração residuos 1..N",
    "",
    "Limitação desta rodada:",
    "  Todos os 330 backbones -> 20 aa",
    "  (comprimento variavel em rodadas futuras)",
]
for i, it in enumerate(items_rfd):
    c = C_VERDE if "  " == it[:2] and "  " == it[2:4] else (C_AMARELO if "Limitação" in it else C_CINZA_CLARO)
    if it.startswith("  "):
        c = C_VERDE
    if "Limitação" in it or "Todos os" in it or "(comprimento" in it:
        c = C_AMARELO
    add_text(slide, it, 0.6, 2.15 + i * 0.33, 5.4, 0.33, font_size=12, color=c)

add_rect(slide, 6.6, 1.25, 6.3, 5.8, C_BRANCO)
add_text(slide, "Backbones gerados (modo REAL)", 6.8, 1.3, 6.0, 0.45,
         font_size=16, bold=True, color=C_AZUL_ESCURO)

comp_rows = [
    ("len5 (cat.)",  "55", "Real — 20 aa"),
    ("len7 (cat.)",  "55", "Real — 20 aa"),
    ("len10 (cat.)", "55", "Real — 20 aa"),
    ("len12 (cat.)", "55", "Real — 20 aa"),
    ("len15 (cat.)", "55", "Real — 20 aa"),
    ("len20 (cat.)", "55", "Real — 20 aa"),
]
hdrs_b = ["Categoria", "Backbones", "Comprimento real"]
ws_b   = [2.1, 1.4, 2.7]
x0b, y0b = 6.7, 1.8
for j, (h, w) in enumerate(zip(hdrs_b, ws_b)):
    add_rect(slide, x0b + sum(ws_b[:j]), y0b, w, 0.4, C_AZUL_MEDIO)
    add_text(slide, h, x0b + sum(ws_b[:j]) + 0.05, y0b + 0.02, w - 0.1, 0.36,
             font_size=12, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
for i, row in enumerate(comp_rows):
    bg = C_CINZA_CLARO if i % 2 == 0 else C_BRANCO
    for j, (cell, w) in enumerate(zip(row, ws_b)):
        add_rect(slide, x0b + sum(ws_b[:j]), y0b + 0.4 + i * 0.38, w, 0.38, bg)
        add_text(slide, cell, x0b + sum(ws_b[:j]) + 0.05, y0b + 0.42 + i * 0.38, w - 0.1, 0.34,
                 font_size=12, color=C_CINZA_TEXTO, align=PP_ALIGN.CENTER)

add_rect(slide, x0b, y0b + 0.4 + 6 * 0.38, sum(ws_b), 0.42, C_AZUL_ESCURO)
add_text(slide, "Total", x0b + 0.05, y0b + 0.42 + 6 * 0.38, ws_b[0] - 0.1, 0.38,
         font_size=13, bold=True, color=C_BRANCO)
add_text(slide, "330", x0b + ws_b[0] + 0.05, y0b + 0.42 + 6 * 0.38, ws_b[1] - 0.1, 0.38,
         font_size=13, bold=True, color=C_VERDE, align=PP_ALIGN.CENTER)
add_text(slide, "todos 20 aa", x0b + sum(ws_b[:2]) + 0.05, y0b + 0.42 + 6 * 0.38, ws_b[2] - 0.1, 0.38,
         font_size=13, bold=True, color=C_AMARELO, align=PP_ALIGN.CENTER)

add_rect(slide, 6.6, 6.35, 6.3, 0.62, RGBColor(0xFF, 0xF3, 0xCD))
add_text(slide,
    "330 backbones conformacionalmente diversos (helices, estendidos, hairpins)\n"
    "Proxima rodada: comprimentos variados independentes (5-15 aa) para sintese",
    6.75, 6.4, 6.1, 0.55, font_size=12, color=C_CINZA_TEXTO)

slide_num(slide, 6)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ProteinMPNN: 24.513 binders
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Design de Sequências — ProteinMPNN REAL",
           "330 backbones × 500 seqs = 165.000 brutas → 24.513 binders únicos 20 aa")

estrategias = [
    ("ProteinMPNN REAL",    "Redesenho da cadeia B (binder 20 aa) para cada backbone",   "modo real", C_VERDE),
    ("Temperatura",         "sampling_temp = 0.1 | backbone_noise = 0.05",               "config",    C_AZUL_MEDIO),
    ("Excluídos",           "Cys, X (--omit_AAs CX) — evitar pontes dissulfeto",          "filtro",    C_AZUL_MEDIO),
    ("BUG CRÍTICO CORRIGIDO","replace('/','') concatenava receptor+binder (240+ aa)",    "Bug 33",    C_VERMELHO),
    ("Fix (commit e9024bc)", "parts = seq.split('/'); binder = parts[-1]  -> 20 aa",     "corrigido", C_VERDE),
    ("Deduplicação",        "165.000 brutas -> 24.513 únicas (83.8% redundância)",        "resultado", C_CIANO),
    ("Homogeneidade",       "Todos os 24.513 binders têm 20 aa (limitação desta rodada)", "nota",      C_AMARELO),
    ("Cache inteligente",   "Não re-roda ProteinMPNN se FASTAs já existem",              "otimizacao",C_AZUL_MEDIO),
    ("is_known_inhibitor",  "Seeds BPTI/SKTI incluídos como âncoras positivas ML",       "ML label",  C_AZUL_MEDIO),
    ("Dataset exportado",   "outputs/dataset/ml_training_dataset.csv — 41 features",    "CSV",       C_VERDE),
]

for i, (nome, desc, tipo, c) in enumerate(estrategias):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.45
    y = 1.25 + row * 1.12
    add_rect(slide, x, y, 6.2, 1.0, C_BRANCO)
    add_rect(slide, x, y, 0.08, 1.0, c)
    add_rect(slide, x + 5.4, y + 0.05, 0.75, 0.38, c)
    add_text(slide, tipo, x + 5.4, y + 0.05, 0.75, 0.38,
             font_size=11, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, nome, x + 0.2, y + 0.06, 5.1, 0.38,
             font_size=13, bold=True, color=C_AZUL_ESCURO)
    add_text(slide, desc, x + 0.2, y + 0.48, 5.1, 0.45,
             font_size=11, color=C_CINZA_TEXTO)

add_rect(slide, 0.4, 6.7, 12.5, 0.65, C_AZUL_ESCURO)
add_text(slide,
    "24.513 binders únicos 20 aa  |  165.000 brutas  |  41 features  |  "
    "Bug FASTA corrigido (commit e9024bc)  |  Cys excluídas",
    0.55, 6.78, 12.1, 0.5, font_size=13, bold=True, color=C_CIANO, align=PP_ALIGN.CENTER)

slide_num(slide, 7)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Dataset ML/DL
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Dataset ML/DL", "24.513 sequências × 41 features + labels Vina/Rosetta reais")

add_text(slide, "Arquivo: outputs/dataset/ml_training_dataset.csv", 0.5, 1.2, 12.0, 0.38,
         font_size=14, color=C_AZUL_MEDIO, bold=True)

grupos = [
    ("Globais (7)", [
        "length, mw_da, net_charge (pH 7)",
        "isoelectric_point, instability_index",
        "aliphatic_index, boman_index",
    ], 0.4, 1.7, C_AZUL_MEDIO),
    ("Hidrofobicidade (2)", [
        "hydrophobicity_kd (Kyte-Doolittle)",
        "frac_hydrophobic (AILMFWV)",
    ], 0.4, 3.5, C_VERDE),
    ("Composição por classe (3)", [
        "frac_aromatic, frac_charged",
        "n_arg_lys",
    ], 0.4, 4.75, C_AZUL_MEDIO),
    ("Features binárias (3)", [
        "has_aromatic_cterminal",
        "has_charged_nterminal, has_pro",
    ], 0.4, 5.8, C_LARANJA),
    ("Composição AA (19)", [
        "frac_A, frac_D, frac_E, frac_F,",
        "frac_G, frac_H, frac_I, frac_K,",
        "frac_L, frac_M, frac_N, frac_P,",
        "frac_Q, frac_R, frac_S, frac_T,",
        "frac_V, frac_W, frac_Y",
    ], 6.5, 1.7, C_AZUL_MEDIO),
    ("Labels supervisionados (3)", [
        "best_affinity_kcal  <- Vina real (194 labels)",
        "rosetta_I_sc  <- PyRosetta I_sc REF2015 (10 labels)",
        "is_known_inhibitor (BPTI/SKTI=1)",
    ], 6.5, 4.5, C_CIANO),
    ("Metadados (2)", [
        "sequence (20 aa), length (categoria)",
        "backbone (stem do PDB RFdiffusion)",
    ], 6.5, 6.1, C_AMARELO),
]

for (titulo, items, x, y, c) in grupos:
    h_box = 0.45 + len(items) * 0.38
    add_rect(slide, x, y, 5.9, h_box, C_BRANCO)
    add_rect(slide, x, y, 5.9, 0.4, c)
    add_text(slide, titulo, x + 0.1, y + 0.03, 5.7, 0.36,
             font_size=13, bold=True, color=C_BRANCO)
    for k, item in enumerate(items):
        add_text(slide, item, x + 0.15, y + 0.45 + k * 0.36, 5.65, 0.34,
                 font_size=11, color=C_CINZA_TEXTO)

slide_num(slide, 8)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Docking: histórico de bugs e resultado
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Docking Molecular — AutoDock Vina",
           "194/194 poses validas | binders RFdiffusion+ProteinMPNN reais | 39 bugs total")

add_rect(slide, 0.4, 1.25, 12.5, 0.55, RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, 0.4, 1.25, 0.1, 0.55, C_VERDE)
add_text(slide,
    "Rodada com binders reais (2026-06-17): 194/194 poses validas | top-1: GSRASARAYAARVRARRAAL (-13.62 kcal/mol)",
    0.6, 1.32, 12.1, 0.42, font_size=13, bold=True, color=C_CINZA_TEXTO)

add_text(slide, "Bugs criticos desta sessao (binders RFdiffusion):", 0.4, 1.9, 12.0, 0.35,
         font_size=14, bold=True, color=C_AZUL_ESCURO)

bugs = [
    ("29-31", "NA atom type rejeitado pelo Vina",
     '"N":"NA" (sodio) em vez de nitrogênio backbone -> "N":"N"', "603cc04", C_VERDE),
    ("32",    "obabel -xr gerava formato de RECEPTOR",
     "Sem ROOT/ENDROOT para ligante -> remover -xr; _ensure_ligand_pdbqt_format() reescreve sempre", "603cc04", C_VERDE),
    ("33 ★",  "TOP CANDIDATOS ERAM REDESIGNS DAS TRIPSINAS (bug critico)",
     "replace('/','') concatenava receptor+binder (240+ aa) -> parts[-1] extrai binder 20 aa", "e9024bc", C_VERMELHO),
    ("34",    "seq[-L:] cortava binders reais de 20 aa",
     "len(seq)>L*2 triggava para len5 (20>10) -> remocao do corte artificial", "ad6fa9e", C_VERDE),
    ("35",    "Grid dimensionado por label (5/7/10), nao comprimento real",
     "item['length']=categoria -> len(item['sequence'])=20 -> grid 80x80x80 A", "ad6fa9e", C_VERDE),
    ("38-39", "Rosetta selecionava por heuristica, nao por Vina | campo vina_affinity inexistente",
     "_select_top_sequences() le docking_results.json | campo correto: best_affinity_kcal", "a7ac84f+e16a28f", C_VERDE),
]
for i, (num, bug, fix, commit, c) in enumerate(bugs):
    y = 2.32 + i * 0.84
    add_rect(slide, 0.4, y, 0.6, 0.72, c)
    add_text(slide, num, 0.4, y + 0.16, 0.6, 0.4,
             font_size=11, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.0, y, 11.9, 0.72, C_BRANCO if i % 2 == 0 else C_CINZA_CLARO)
    add_text(slide, bug, 1.1, y + 0.02, 7.5, 0.32,
             font_size=12, bold=True, color=C_AZUL_ESCURO if c != C_VERMELHO else C_VERMELHO)
    add_text(slide, fix, 1.1, y + 0.35, 7.5, 0.32,
             font_size=10, color=C_CINZA_TEXTO)
    add_text(slide, commit, 8.65, y + 0.18, 2.2, 0.32,
             font_size=10, color=C_AZUL_MEDIO, align=PP_ALIGN.CENTER)

slide_num(slide, 9)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Resultados docking com binders reais
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Resultados Docking + Ranking — Rodada Expandida (2026-06-27)",
           "880/1000 poses validas | novo top-1: SARESIKKAYKTFLERYKKL -14,58 kcal/mol | 24.513 rankeados")

add_rect(slide, 0.4, 1.25, 12.5, 1.0, C_AZUL_ESCURO)
add_text(slide, "880 / 1000", 0.5, 1.28, 4.5, 0.65,
         font_size=44, bold=True, color=C_VERDE, align=PP_ALIGN.CENTER)
add_text(slide, "poses Vina validas (rodada expandida)", 0.5, 1.88, 4.5, 0.28,
         font_size=13, color=C_CINZA_CLARO, align=PP_ALIGN.CENTER)
stats_r = [
    ("Melhor:", "-14.58 kcal/mol", C_VERDE),
    ("Media:", "~-12.4 kcal/mol", C_CIANO),
    ("ML RMSE:", "0.514 kcal/mol", C_AMARELO),
]
for k, (lbl, val, c) in enumerate(stats_r):
    add_text(slide, lbl, 5.2 + k * 2.5, 1.35, 1.1, 0.35,
             font_size=12, color=C_CINZA_CLARO)
    add_text(slide, val, 5.2 + k * 2.5, 1.66, 2.2, 0.42,
             font_size=18, bold=True, color=c)

add_text(slide, "Top-5 por score composto (ranking expandido):", 0.4, 2.35, 8.0, 0.38,
         font_size=14, bold=True, color=C_AZUL_ESCURO)

top_r = [
    (1, "SARESIKKAYKTFLERYKKL", 20, -14.58, 0.748, "NOVO top-1 — Vina mais negativo do pipeline"),
    (2, "GSRASARAYAARVRARRAAL", 20, -13.62, 0.701, "Validado por MD (RMSD 0,494 nm, ESTAVEL)"),
    (3, "GARESIREHQKRFLERYKKK", 20, -13.56, 0.661, "Gly/Ala + Glu/His + R/K C-term"),
    (4, "GGPTGKRIAELYKKSLEKKK", 20, -13.47, 0.653, "Gly/Pro scaffold + Lys-rico"),
    (5, "AARENIRAYAARFRARLAAK", 20, -13.79, 0.635, "Ala-rico + Arg/Tyr — Vina entre melhores"),
]
hdrs_r = ["#", "Sequencia (20 aa)", "aa", "Vina", "Score", "Observacao"]
ws_r   = [0.38, 4.0, 0.42, 1.2, 0.85, 5.47]
x0r, y0r = 0.4, 2.8

for j, (h, w) in enumerate(zip(hdrs_r, ws_r)):
    add_rect(slide, x0r + sum(ws_r[:j]), y0r, w, 0.35, C_AZUL_ESCURO)
    add_text(slide, h, x0r + sum(ws_r[:j]) + 0.04, y0r + 0.02, w - 0.08, 0.31,
             font_size=10, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)

for i, (rank, seq, aa, vina, score, obs) in enumerate(top_r):
    bg = C_CINZA_CLARO if i % 2 == 0 else C_BRANCO
    add_rect(slide, x0r, y0r + 0.35 + i * 0.38, 0.07, 0.38, C_VERDE if rank <= 2 else C_AZUL_MEDIO)
    row_vals = [str(rank), seq, str(aa), f"{vina:.2f}", f"{score:.3f}", obs]
    for j, (cell, w) in enumerate(zip(row_vals, ws_r)):
        add_rect(slide, x0r + sum(ws_r[:j]), y0r + 0.35 + i * 0.38, w, 0.38, bg)
        c_cell = C_VERDE if j == 3 and rank == 1 else (C_AZUL_ESCURO if j == 1 else C_CINZA_TEXTO)
        add_text(slide, cell, x0r + sum(ws_r[:j]) + 0.04, y0r + 0.37 + i * 0.38, w - 0.08, 0.34,
                 font_size=9 if j in (1,5) else 10, bold=(j==1 and rank<=2),
                 color=c_cell, align=PP_ALIGN.LEFT if j in (1,5) else PP_ALIGN.CENTER)

add_rect(slide, 0.4, 7.1, 12.5, 0.28, C_AZUL_MEDIO)
add_text(slide,
    "Especificidade: 20/20 aprovados vs tripsina humana (1TRN) + Apis mellifera (A0A7M7MMI1) | "
    "Resistencia: 0/20 resistentes — todos susceptiveis (K/R internos) — sondas in vitro",
    0.55, 7.12, 12.2, 0.24, font_size=11, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)

slide_num(slide, 10)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — PyRosetta: I_sc × Vina cross-validation
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "PyRosetta 2026.25 — Refinamento de Interface",
           "FastRelax + InterfaceAnalyzerMover('A_B') | I_sc REF2015 | 10/10 complexos refinados")

add_rect(slide, 0.4, 1.25, 12.5, 0.65, RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, 0.4, 1.25, 0.1, 0.65, C_VERDE)
add_text(slide,
    "Instalado: pyrosetta-2026.25+release (pyrosetta-installer, sem licenca separada) | "
    "Python 3.10 | protein_design_env | 2026-06-17",
    0.6, 1.35, 12.1, 0.48, font_size=13, bold=True, color=C_CINZA_TEXTO)

# Tabela Vina x Rosetta
add_text(slide, "Cross-validation Vina × I_sc Rosetta — top-10 candidatos:", 0.4, 2.0, 12.0, 0.38,
         font_size=14, bold=True, color=C_AZUL_ESCURO)

cross = [
    ("GARKSIREYQKRVLERLKKK", -12.76, 2, -86.28, 1, "+1", "bestI_sc"),
    ("AARASQREYQKKFLERLKKK", -12.52, 5, -85.92, 2, "+3", "elevado"),
    ("MKKQRENAKKVAEITLKKAK", -12.68, 4, -80.49, 3, "+1", "concordante"),
    ("GSRASARAYAARVRARRAAL", -13.62, 1, -78.44, 4, "-3", "bestVina"),
    ("AARASIRAAAARFRARRAAL", -12.62, 6, -75.45, 5, "+1", "concordante"),
    ("GSLTGRRIAALWKASLAKRK", None,   "-",-69.44, 6, "-",  "Rosetta-only"),
    ("SAAARARQRAVGARMRARVA", -12.44, 8, -63.98, 7, "+1", "concordante"),
    ("AARENIRKAHKTFLERLKKK", -12.36, 7, -62.74, 8, "-1", "concordante"),
    ("ALDAVRARARALGARLRARA", None,   "-",-61.89, 9, "-",  "Rosetta-only"),
    ("SLARKRAEENAKRFLERVKK", -12.71, 3, -61.41, 10,"-7", "divergente"),
]

hdrs_c = ["Sequencia", "Vina", "RankV", "I_sc", "RankR", "Delta", "Status"]
ws_c   = [4.2, 1.2, 0.8, 1.4, 0.8, 0.8, 2.0]
x0c, y0c = 0.4, 2.45

for j, (h, w) in enumerate(zip(hdrs_c, ws_c)):
    add_rect(slide, x0c + sum(ws_c[:j]), y0c, w, 0.35, C_AZUL_ESCURO)
    add_text(slide, h, x0c + sum(ws_c[:j]) + 0.04, y0c + 0.02, w - 0.08, 0.31,
             font_size=10, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)

status_colors = {
    "bestI_sc": C_VERDE, "elevado": C_VERDE, "concordante": C_AZUL_MEDIO,
    "bestVina": C_CIANO, "Rosetta-only": C_CINZA_TEXTO, "divergente": C_VERMELHO,
}
for i, (seq, vina, rankv, isc, rankr, delta, status) in enumerate(cross):
    bg = C_CINZA_CLARO if i % 2 == 0 else C_BRANCO
    sc = status_colors.get(status, C_CINZA_TEXTO)
    vina_str = f"{vina:.2f}" if vina else "-"
    row_vals = [seq, vina_str, str(rankv), f"{isc:.2f}", str(rankr), delta, status]
    for j, (cell, w) in enumerate(zip(row_vals, ws_c)):
        add_rect(slide, x0c + sum(ws_c[:j]), y0c + 0.35 + i * 0.30, w, 0.30, bg)
        cell_c = sc if j == 6 else (C_CINZA_TEXTO if j != 1 else (C_VERDE if vina and float(vina) <= -13.0 else C_CINZA_TEXTO))
        if j == 3 and float(isc) <= -80:
            cell_c = C_VERDE
        add_text(slide, cell, x0c + sum(ws_c[:j]) + 0.04, y0c + 0.37 + i * 0.30, w - 0.08, 0.26,
                 font_size=9 if j in (0,6) else 10,
                 color=cell_c, align=PP_ALIGN.LEFT if j in (0,6) else PP_ALIGN.CENTER)

add_rect(slide, 0.4, 5.55, 12.5, 1.7, C_BRANCO)
add_rect(slide, 0.4, 5.55, 0.06, 1.7, C_VERDE)
insights = [
    "GARKSIREYQKRVLERLKKK: melhor I_sc (-86.28) apesar de #2 por Vina — FastRelax revela interface forte",
    "AARASQREYQKKFLERLKKK: sobe de #5 Vina para #2 I_sc (-85.92) — docking rigido subestimava sua afinidade",
    "SLARKRAEENAKRFLERVKK: cai de #3 Vina para #10 I_sc (-61.41) — DESCARTADO, interface nao sustenta pos-relax",
    "GSRASARAYAARVRARRAAL: confirmado por ambos (Vina #1, I_sc #4) — candidato mais robusto para sintese",
    "Concordancia Vina × Rosetta: criterio mais rigoroso — identifica GARKSIREYQKRVLERLKKK como top real",
]
for i, ins in enumerate(insights):
    c_ins = C_VERDE if "GARKS" in ins or "GSRASA" in ins else (C_VERMELHO if "DESCARTADO" in ins else C_CINZA_TEXTO)
    add_text(slide, ins, 0.6, 5.62 + i * 0.32, 12.1, 0.30, font_size=11, color=c_ins)

slide_num(slide, 11)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Top candidatos prioritários (Vina+Rosetta)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Candidatos Prioritarios — Top-3 Confirmados por Dupla Validacao",
           "Vina (docking rigido) × PyRosetta I_sc (FastRelax) | todos 20 aa | Fmoc-SPPS")

top_cands = [
    ("GARKSIREYQKRVLERLKKK", -12.76, -86.28, 2, 1,
     "R/K denso (n=9), Glu/Tyr distribuido",
     "Melhor I_sc absoluto. FastRelax revela interface excepcional. Melhor candidato por concordancia.", C_VERDE),
    ("GSRASARAYAARVRARRAAL", -13.62, -78.44, 1, 4,
     "Gly/Ser N-term, Arg/Ala backbone regular",
     "Melhor Vina absoluto (-13.62). Confirmado por Rosetta (#4 I_sc). Candidato de referencia.", C_CIANO),
    ("AARASQREYQKKFLERLKKK", -12.52, -85.92, 5, 2,
     "Ala N-term, Arg/Lys/Phe C-term",
     "Vina #5 mas Rosetta #2. Docking rigido subestimava afinidade real de interface.", C_VERDE),
    ("MKKQRENAKKVAEITLKKAK", -12.68, -80.49, 4, 3,
     "Lys-rico (n=8), alifatico, solúvel",
     "Consistente em ambos (#4 Vina, #3 Rosetta). Alto conteudo Lys favorece sintese e solubilidade.", C_AZUL_MEDIO),
    ("SLARKRAEENAKRFLERVKK", -12.71, -61.41, 3, 10,
     "Leu/Ala scaffold + R/K alternados",
     "DESCARTADO: Vina #3 mas Rosetta #10. Interface nao sustenta apos FastRelax. Nao sintetizar.", C_VERMELHO),
]

for i, (seq, vina, isc, rv, ri, perfil, desc, c) in enumerate(top_cands):
    y = 1.28 + i * 1.2
    descartado = (c == C_VERMELHO)
    bg_card = RGBColor(0xFF, 0xEE, 0xEE) if descartado else C_BRANCO
    add_rect(slide, 0.4, y, 12.5, 1.1, bg_card)
    add_rect(slide, 0.4, y, 0.45, 1.1, c)
    prioridade = "✗" if descartado else f"#{i+1}"
    add_text(slide, prioridade, 0.4, y + 0.28, 0.45, 0.5,
             font_size=20, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.95, y + 0.06, 2.2, 0.28, c)
    add_text(slide, perfil[:30], 0.97, y + 0.07, 2.18, 0.26,
             font_size=10, color=C_BRANCO, align=PP_ALIGN.CENTER)

    add_text(slide, seq, 3.25, y + 0.03, 5.2, 0.38,
             font_size=14, bold=True, color=(C_VERMELHO if descartado else C_AZUL_ESCURO))
    add_text(slide, f"Vina = {vina:.2f} kcal/mol (rank #{rv})  |  I_sc = {isc:.2f} kcal/mol (rank #{ri})", 3.25, y + 0.42, 5.2, 0.28,
             font_size=11, color=C_AZUL_MEDIO)
    add_text(slide, desc, 3.25, y + 0.70, 5.2, 0.36,
             font_size=10, italic=True, color=C_CINZA_TEXTO)

    bv = abs(vina - 10.0) / 4.0
    br = abs(isc - 50.0) / 40.0
    add_rect(slide, 8.55, y + 0.18, bv * 1.8, 0.22, C_CIANO)
    add_text(slide, f"V:{vina:.1f}", 8.55 + bv * 1.8 + 0.05, y + 0.18, 0.7, 0.22,
             font_size=9, color=C_CIANO)
    add_rect(slide, 8.55, y + 0.50, br * 1.8, 0.22, C_VERDE)
    add_text(slide, f"R:{isc:.0f}", 8.55 + br * 1.8 + 0.05, y + 0.50, 0.7, 0.22,
             font_size=9, color=C_VERDE)

add_rect(slide, 0.4, 7.22, 12.5, 0.18, C_VERDE)
add_text(slide,
    "MD 11/11 CONCLUIDO (Fases 3+4): top-3 sintese — MKKQRENA(0,447nm) | RLREELKK(0,294nm★) | GSRASARA(0,494nm) | ver slide 13",
    0.55, 7.23, 12.2, 0.16, font_size=11, color=C_BRANCO, align=PP_ALIGN.CENTER)

slide_num(slide, 12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — MD Results: 5 candidatos × ACR157 (10 ns cada)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Dinamica Molecular — 11 Candidatos × ACR157 (Fases 3+4)",
           "GROMACS | AMBER99SB-ILDN | TIP3P | 300 K | 10 ns cada | RTX 5070 Ti (~36 min/replica)")

# — Cabeçalho tabela
add_rect(slide, 0.3, 1.25, 12.7, 0.42, C_AZUL_ESCURO)
for txt, x, w in [("Candidato", 0.35, 3.55), ("Fase", 4.0, 0.55), ("Vina", 4.6, 1.1),
                   ("I_sc", 5.7, 1.1), ("RMSD avg±DP (nm)", 6.85, 2.8), ("Rg", 9.7, 0.8), ("Estabilidade", 9.75, 3.2)]:
    add_text(slide, txt, x, 1.3, w, 0.32, font_size=11, bold=True, color=C_BRANCO)

# — Linhas da tabela
C_VERDE_MD   = RGBColor(0xD4, 0xED, 0xDA)
C_AMARELO_MD = RGBColor(0xFF, 0xF3, 0xCD)
C_VERMELHO_MD= RGBColor(0xF8, 0xD7, 0xDA)

md_rows = [
    # Fase 3
    ("MKKQRENAKKVAEITLKKAK", "-12,72", "-80,49", "0,447 ± 0,332", "1,785", "F3", "ESTAVEL",   C_VERDE_MD,    C_VERDE),
    ("GSRASARAYAARVRARRAAL", "-13,62", "-78,44", "0,494 ± 0,414", "1,799", "F3", "ESTAVEL",   C_VERDE_MD,    C_VERDE),
    ("AARASIRAAAARFRARRAAL", "-12,62", "-75,45", "0,945 ± 0,924", "1,836", "F3", "Marginal",  C_AMARELO_MD,  C_AMARELO),
    # Fase 4
    ("RLREELKKAEEWLEKRRKEE", "—",      "—",      "0,294 ± 0,065", "1,780", "F4", "ESTAVEL*",  C_VERDE_MD,    C_VERDE),
    ("RLRAIWLEAEKLLEERRKKK", "—",      "—",      "0,725 ± 0,870", "1,815", "F4", "Marginal",  C_AMARELO_MD,  C_AMARELO),
    ("SARESIKKAYKTFLERYKKL", "-14,58", "—",      "0,871 ± 0,844", "1,864", "F4", "Marginal",  C_AMARELO_MD,  C_AMARELO),
    ("KRLRENWLEAEKLLEERRKKK","—",      "—",      "1,074 ± 1,103", "1,869", "F4", "Instavel",  C_VERMELHO_MD, C_LARANJA),
]
for i, (seq, vina, isc, rmsd, rg, fase, status, bg, sc) in enumerate(md_rows):
    y = 1.72 + i * 0.67
    add_rect(slide, 0.3, y, 12.7, 0.62, bg)
    add_rect(slide, 0.3, y, 0.08, 0.62, sc)
    add_text(slide, seq, 0.45, y + 0.14, 3.55, 0.32, font_size=10, bold=True, color=C_AZUL_ESCURO)
    add_text(slide, fase, 4.0, y + 0.14, 0.55, 0.32, font_size=9, bold=True,
             color=C_VERDE if fase=="F4" else C_AZUL_MEDIO)
    for val, x in [(vina, 4.6), (isc, 5.7), (rmsd, 6.85), (rg, 8.55)]:
        add_text(slide, val, x, y + 0.14, 1.15, 0.32, font_size=10, color=C_CINZA_TEXTO)
    add_rect(slide, 9.75, y + 0.1, 3.2, 0.42, sc)
    add_text(slide, status, 9.75, y + 0.1, 3.2, 0.42,
             font_size=11, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)

# — Painel inferior: reclassificação
add_rect(slide, 0.3, 6.44, 12.7, 0.88, C_AZUL_ESCURO)
add_text(slide, "TOP-3 SINTESE (Fases 3+4) — todos RMSD < 0,5 nm em 10 ns", 0.5, 6.48, 12.3, 0.32,
         font_size=12, bold=True, color=C_CIANO)
add_text(slide,
    "#1 MKKQRENAKKVAEITLKKAK (0,447 nm, F3)  |  "
    "#2 RLREELKKAEEWLEKRRKEE (0,294 nm, F4 — MAIS ESTAVEL DO PIPELINE)  |  "
    "#3 GSRASARAYAARVRARRAAL (0,494 nm, F3)  |  "
    "SARESIKKAYKTFLERYKKL: melhor Vina (-14,58) mas marginal em MD",
    0.5, 6.82, 12.3, 0.46, font_size=10, color=C_CINZA_CLARO)

slide_num(slide, 13)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Ferramentas: status atual no servidor
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Status das Ferramentas no Servidor",
           "eulalio@200.235.143.10 | Debian | RTX 5070 Ti 16 GB | sm_120 Blackwell | 2026-06-18")

ferramentas = [
    ("✓", "Python 3.10",               "conda env protein_design_env (pipeline principal)",            C_VERDE),
    ("✓", "OpenBabel",                 "obabel — conversao PDBQT (sem -xr para ligantes)",             C_VERDE),
    ("✓", "AutoDock Vina f458505-mod", "docking rigido peptideos (TORSDOF=0, best_affinity_kcal)",    C_VERDE),
    ("✓", "fpocket",                   "analise e deteccao de bolso de ligacao",                       C_VERDE),
    ("✓", "GROMACS gmx_mpi",           "CUDA-MPI, build sm_120 Blackwell, env md-gromacs",            C_VERDE),
    ("✓", "PyTorch 2.12.0.dev+cu128", "CUDA 12.8, compativel Blackwell (sm_120)",                    C_VERDE),
    ("✓", "ProteinMPNN",               "~/ProteinMPNN | 24.513 binders gerados para ACR157",          C_VERDE),
    ("✓", "RFdiffusion",               "~/RFdiffusion | Complex_base_ckpt.pt 462 MB | 330 backbones", C_VERDE),
    ("✓", "PyRosetta 2026.25",         "pyrosetta-installer (sem licenca) | FastRelax + IA Mover",    C_VERDE),
    ("✓", "plip 3.0.0",                "analise de interacoes protein-ligante (conda-forge)",          C_VERDE),
    ("✓", "rdkit, pdbfixer, mdtraj",   "analise estrutural e preparacao de sistemas MD",               C_VERDE),
    ("✓", "PeptideBuilder",            "geracao all-atom de peptideos (fallback + Rosetta complex)",   C_VERDE),
]

for i, (status, nome, desc, c) in enumerate(ferramentas):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.45
    y = 1.25 + row * 0.97
    add_rect(slide, x, y, 6.2, 0.88, C_BRANCO if row % 2 == 0 else C_CINZA_CLARO)
    add_rect(slide, x, y, 0.55, 0.88, c)
    add_text(slide, status, x, y + 0.2, 0.55, 0.45,
             font_size=20, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, nome, x + 0.65, y + 0.06, 5.45, 0.38,
             font_size=14, bold=True, color=C_AZUL_ESCURO)
    add_text(slide, desc, x + 0.65, y + 0.46, 5.45, 0.38,
             font_size=11, color=C_CINZA_TEXTO)

slide_num(slide, 14)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Estratégia de Docking Rígido  (renumerado +1)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Estrategia Tecnica: Docking Rigido para ML Labels",
           "Justificativa cientifica e tecnica | binders RFdiffusion 20 aa | TORSDOF 0")

boxes = [
    (0.4, 1.3, 2.8, 1.2, "Binder\nall-atom", "PeptideBuilder\n(ECEPP/3)", C_AZUL_MEDIO),
    (3.6, 1.3, 2.8, 1.2, "PDBQT\nall-atom",  "_pdb_to_pdbqt_minimal()\n(N->N, nao NA!)", C_AZUL_MEDIO),
    (6.8, 1.3, 2.8, 1.2, "ROOT wrapper",     "_ensure_ligand\n_pdbqt_format()", C_VERDE),
    (10.1,1.3, 2.8, 1.2, "AutoDock\nVina",   "Grid 80x80x80 A\nexh=8, modes=9", C_LARANJA),
]
for (x, y, w, h, t1, t2, c) in boxes:
    add_rect(slide, x, y, w, h, c)
    add_text(slide, t1, x + 0.1, y + 0.07, w - 0.2, 0.5,
             font_size=15, bold=True, color=C_BRANCO)
    add_text(slide, t2, x + 0.1, y + 0.6, w - 0.2, 0.55,
             font_size=12, color=C_CINZA_CLARO)
for x in [3.2, 6.4, 9.65]:
    add_text(slide, "->", x + 0.1, 1.7, 0.5, 0.5,
             font_size=28, bold=True, color=C_CIANO, align=PP_ALIGN.CENTER)

add_rect(slide, 0.4, 2.7, 12.5, 0.08, C_CIANO)
add_text(slide, "Por que docking RIGIDO?", 0.4, 2.85, 6.0, 0.42,
         font_size=18, bold=True, color=C_AZUL_ESCURO)

razoes = [
    ("Limite tecnico:", "Peptideos 20 aa tem >32 ligacoes rotacionaveis (phi+psi+chi). "
     "Vina aborta docking flexivel acima de 32 torsional bonds."),
    ("Escala:", "Para triagem de 24.513 binders (ML labels), docking rigido e viavel. "
     "Correlacao Pearson ~0.7 com flexivel para top candidatos."),
    ("Custo:", "194 candidatos x 8 exhaustiveness ~ 15 min. "
     "Docking flexivel completo seria inviavel nessa escala."),
    ("Validacao posterior:", "Top-10 refinados com PyRosetta FastRelax (I_sc REF2015) "
     "e top-5 por MD GROMACS 10 ns — validacao rigorosa em camadas."),
]
for i, (titulo, texto) in enumerate(razoes):
    add_rect(slide, 0.4, 3.35 + i * 0.95, 12.5, 0.88, C_BRANCO if i % 2 == 0 else C_CINZA_CLARO)
    add_text(slide, titulo, 0.6, 3.38 + i * 0.95, 3.5, 0.38,
             font_size=13, bold=True, color=C_AZUL_MEDIO)
    add_text(slide, texto, 4.1, 3.38 + i * 0.95, 8.7, 0.85,
             font_size=13, color=C_CINZA_TEXTO)

slide_num(slide, 15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Infraestrutura e Workflow
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Infraestrutura Computacional e Workflow Git",
           "Notebook Windows 11 <-> GitHub <-> Servidor Debian RTX 5070 Ti")

nos = [
    (0.4,  2.8, 3.8, 2.0, "Notebook Local",
     "Windows 11\nRTX 2050 4 GB\nEdicao de codigo", C_AZUL_MEDIO),
    (4.7,  2.8, 3.8, 2.0, "GitHub",
     "eulaliobqi/design-inibidores\ngit push / pull\nFonte de verdade", C_CINZA_TEXTO),
    (9.0,  2.8, 3.8, 2.0, "Servidor",
     "eulalio@200.235.143.10\nDebian | 32 cores\nRTX 5070 Ti 16 GB", C_VERDE),
]
for (x, y, w, h, t, d, c) in nos:
    add_rect(slide, x, y, w, h, c)
    add_text(slide, t, x + 0.15, y + 0.12, w - 0.3, 0.5,
             font_size=14, bold=True, color=C_BRANCO)
    add_text(slide, d, x + 0.15, y + 0.65, w - 0.3, 1.2,
             font_size=12, color=C_CINZA_CLARO)

setas = [(4.3, 3.75, "git push"), (8.6, 3.75, "git pull")]
for (x, y, label) in setas:
    add_text(slide, "->", x, y - 0.05, 0.5, 0.5,
             font_size=24, bold=True, color=C_CIANO, align=PP_ALIGN.CENTER)
    add_text(slide, label, x - 0.15, y + 0.4, 0.8, 0.3,
             font_size=11, color=C_CIANO, align=PP_ALIGN.CENTER)

add_rect(slide, 0.4, 5.1, 12.5, 0.08, C_CIANO)
add_text(slide, "Comandos padrao no servidor:", 0.5, 5.3, 5.0, 0.38,
         font_size=15, bold=True, color=C_AZUL_ESCURO)

cmds = [
    "cd ~/design-inibidores && git pull origin main",
    "conda run -n protein_design_env python scripts/run_pipeline.py --step <step> --resume",
    "conda run -n md-gromacs python scripts/run_pipeline.py --step md --resume   # apenas MD",
    "screen -S <nome>   # OBRIGATORIO antes de qualquer run longa (SIGTTOU mata background)",
]
add_rect(slide, 0.4, 5.75, 12.5, 1.6, C_AZUL_ESCURO)
for i, cmd in enumerate(cmds):
    add_text(slide, "$ " + cmd, 0.55, 5.82 + i * 0.38, 12.2, 0.35,
             font_size=12, color=C_VERDE)

slide_num(slide, 16)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 17 — Próximos Passos
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_light(slide)
header_bar(slide, "Proximos Passos", "Fase 4-5 | resultados Fase 1-3 concluidos (2026-06-27)")

steps_prox = [
    ("4a", "PROXIMO",  "RFdiffusion 5-15 aa + filtro KR-interno=0",
     "250 backbones curtos: 50 x [5, 7, 10, 12, 15 aa] | contigs independentes por comprimento\n"
     "ProteinMPNN: P1 unico C-terminal (Arg/Lys) | posicoes internas = Ala/Ser/Gly/Pro",
     C_VERDE, "~3h GPU"),
    ("4b", "PROXIMO",  "MD dos 219 candidatos OptimizationAgent",
     "219 variantes de MKKQRENAKKVAEITLKKAK + GSRASARAYAARVRARRAAL geradas\n"
     "Selecionar top-5 por Vina -> MD 10 ns -> filtrar estaveis (RMSD < 0,5 nm)",
     C_CIANO, "~5h GPU"),
    ("4c", "PROXIMO",  "MD de SARESIKKAYKTFLERYKKL (novo top-1 Vina)",
     "Melhor afinidade Vina do pipeline (-14,58 kcal/mol) ainda sem MD\n"
     "Confirmar estabilidade antes de priorizar para sintese",
     C_AMARELO, "~1h GPU"),
    ("5",  "FUTURO",   "Expansao taxonomica: H. armigera + A. gemmatalis + D. saccharalis",
     "Pipeline atual: ACR157 (A. gemmatalis) apenas\n"
     "Rodar docking + Rosetta nos outros 3 receptores -> seletividade por especie",
     C_AZUL_MEDIO, "~GPU"),
    ("6a", "FUTURO",   "Variantes resistentes: Nle/Orn/D-aa nos P1-internos",
     "K interno -> Nle (norleucina) | R interno -> Orn (ornitina)\n"
     "D-Lys / D-Arg: invisiveis para L-tripsinas enantiosseletivas (Fmoc-SPPS)",
     C_LARANJA, "sintese"),
    ("6b", "LONGO",    "Sintese e validacao experimental",
     "Top-2 estaveis -> Fmoc-SPPS -> IC50 in vitro (S. frugiperda + A. gemmatalis)\n"
     "Publicacao alvo: JCIM / CSBJ",
     C_LARANJA, "meses"),
]

for i, (num, prazo, titulo, desc, c, tempo) in enumerate(steps_prox):
    col = i % 2
    row = i // 2
    if i == 5:
        x, y, w = 0.4, 1.25 + 2 * 1.85, 12.5
    else:
        x = 0.4 + col * 6.5
        y = 1.25 + row * 1.85
        w = 6.2
    add_rect(slide, x, y, w, 1.65, C_BRANCO)
    add_rect(slide, x, y, 0.55, 1.65, c)
    add_text(slide, num, x, y + 0.45, 0.55, 0.5,
             font_size=22, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 0.6, y + 0.04, 1.3, 0.3, c)
    add_text(slide, prazo, x + 0.6, y + 0.04, 1.3, 0.3,
             font_size=10, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, titulo, x + 0.65, y + 0.38, w - 2.0, 0.38,
             font_size=13, bold=True, color=C_AZUL_ESCURO)
    add_text(slide, desc, x + 0.65, y + 0.76, w - 2.0, 0.85,
             font_size=11, color=C_CINZA_TEXTO)
    add_text(slide, tempo, x + w - 1.5, y + 0.05, 1.4, 0.35,
             font_size=11, bold=True, color=c, align=PP_ALIGN.RIGHT)

slide_num(slide, 17)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 18 — Conclusões
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
bg_dark(slide)
add_rect(slide, 0, 0, 13.33, 1.25, C_AZUL_ESCURO)
add_rect(slide, 0, 1.2, 13.33, 0.08, C_CIANO)
add_text(slide, "Conclusoes", 0.5, 0.15, 12.0, 0.6,
         font_size=34, bold=True, color=C_BRANCO)
add_text(slide, "Fases 1-4 CONCLUIDAS | 11 MD | 3 candidatos estaveis | RLREELKK mais estavel do pipeline — 2026-06-27", 0.5, 0.72, 12.0, 0.42,
         font_size=16, color=C_CIANO)

conclusoes = [
    (C_VERDE, "✓",
     "330 backbones RFdiffusion | 24.513 binders ProteinMPNN | 880 poses Vina | "
     "top-1 Vina: SARESIKKAYKTFLERYKKL (-14,58 kcal/mol) | 219 variantes OptimizationAgent"),
    (C_VERDE, "✓",
     "MD 11/11 concluidos (Fases 3+4) | TOP-3 SINTESE: "
     "#1 MKKQRENAKKVAEITLKKAK (0,447 nm) | #2 RLREELKKAEEWLEKRRKEE (0,294 nm ★MAIS ESTAVEL) | "
     "#3 GSRASARAYAARVRARRAAL (0,494 nm)"),
    (C_VERDE, "✓",
     "Especificidade: 20/20 aprovados vs tripsina humana (1TRN) + Apis mellifera (SI >= 2,0 kcal/mol) | "
     "ML RF RMSE=0,514 kcal/mol | 24.513 predicoes geradas"),
    (C_VERDE, "✓",
     "Resistencia proteolitica: 0/20 resistentes (3-7 P1-internos K/R) — sondas in vitro. "
     "Especificidade preservada: seletividade nao e comprometida pelo perfil Arg/Lys"),
    (C_CIANO, "→",
     "MD obrigatorio como filtro: GARKSIREYQKRVLERLKKK (melhor I_sc -86,28) instavel em MD (1,45 nm). "
     "SARESIKKAYKTFLERYKKL (melhor Vina -14,58) marginal em MD (0,871 nm). "
     "Afinidade estatica != estabilidade dinamica"),
    (C_CIANO, "→",
     "RLREELKKAEEWLEKRRKEE: surpresa positiva — 5o por heuristica OptimizationAgent mas mais estavel do pipeline "
     "(RMSD 0,294 nm, DP=0,065). Composicao Glu/Leu sugere mecanismo de ligacao distinto"),
    (C_LARANJA, "⏳",
     "Proximo: RFdiffusion 5-15 aa (KR-interno=0) + estrategias Nle/Orn/D-aa para resistencia proteolitica | "
     "Sintese: MKKQRENA + RLREELKK + GSRASARA -> Fmoc-SPPS -> IC50 in vitro (JCIM/CSBJ)"),
]
for i, (c, sym, texto) in enumerate(conclusoes):
    add_rect(slide, 0.4, 1.4 + i * 0.82, 0.55, 0.7, c)
    add_text(slide, sym, 0.4, 1.47 + i * 0.82, 0.55, 0.55,
             font_size=20, bold=True, color=C_BRANCO, align=PP_ALIGN.CENTER)
    add_text(slide, texto, 1.05, 1.44 + i * 0.82, 12.0, 0.72,
             font_size=12.5, color=C_CINZA_CLARO)

slide_num(slide, 18)

# ─── Salvar ───────────────────────────────────────────────────────────────────
out = "apresentacao_design_inibidores.pptx"
prs.save(out)
print(f"OK  Salvo: {out}  ({os.path.getsize(out)//1024} KB)")
